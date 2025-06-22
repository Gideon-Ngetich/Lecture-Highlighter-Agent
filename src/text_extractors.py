import os
import docx
import PyPDF2
from pptx import Presentation

def extract_text_from_docx(file_path):
    """Extract text from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX: {str(e)}")

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text.strip()
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF: {str(e)}")

def extract_text_from_ppt(file_path):
    """Extract text from a PPT file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error extracting text from PPT: {str(e)}")

def extract_text(file_path):
    """Determine file type and extract text accordingly."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")